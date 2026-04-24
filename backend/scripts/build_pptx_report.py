#!/usr/bin/env python3
from __future__ import annotations

import argparse
import copy
import io
import json
import math
import re
import sys
import urllib.request
from datetime import datetime
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx_template.src.renderer import build_report


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build PPTX report from template/sample JSON.")
    parser.add_argument("--out", required=True, help="Absolute path to generated .pptx")
    return parser.parse_args()


def safe_num(v: object, ndigits: int = 1) -> str:
    try:
        n = float(v)
        if n != n:
            return "—"
        return str(round(n, ndigits))
    except Exception:
        return "—"


def str_or_default(v: object, fallback: str = "—") -> str:
    if v is None:
        return fallback
    s = str(v).strip()
    return s if s else fallback


def clean_address(raw: object) -> str:
    s = str_or_default(raw, "")
    if not s:
        return "—"
    # Убираем служебный префикс НСПД до фактического адреса региона.
    s = re.sub(
        r"(?is)^.*?адрес\s+ориентир[а-я]*[:\s,-]*",
        "",
        s,
    ).strip(" ,.-")
    return s or "—"


def format_sotok(area_m2: object) -> str:
    try:
        sotok = float(area_m2) / 100.0
    except Exception:
        return "—"
    if sotok != sotok:  # NaN
        return "—"

    rounded = round(sotok, 1)
    has_fraction = not float(rounded).is_integer()
    if has_fraction:
        value = f"{rounded:.1f}".replace(".", ",")
        return f"{value} соток"

    n = int(round(rounded))
    value = str(n)
    mod10 = n % 10
    mod100 = n % 100
    if 11 <= mod100 <= 14:
        noun = "соток"
    elif mod10 == 1:
        noun = "сотка"
    elif 2 <= mod10 <= 4:
        noun = "сотки"
    else:
        noun = "соток"
    return f"{value} {noun}"


def format_check_date_ru(now: datetime) -> str:
    months = {
        1: "января",
        2: "февраля",
        3: "марта",
        4: "апреля",
        5: "мая",
        6: "июня",
        7: "июля",
        8: "августа",
        9: "сентября",
        10: "октября",
        11: "ноября",
        12: "декабря",
    }
    return f"{now.day} {months[now.month]} {now.year}"


def risk_level(score: int) -> tuple[str, str]:
    if score >= 80:
        return "Высокий риск", "high.png"
    if score >= 40:
        return "Средний риск", "medium.png"
    return "Низкий риск", "zero.png"


def terrain_level_name(max_slope: float | None) -> tuple[str, str]:
    if max_slope is None:
        return "Не определен", "zero.png"
    if max_slope > 15:
        return "Критический", "high.png"
    if max_slope > 8:
        return "Средний", "medium.png"
    return "Низкий", "zero.png"


def with_units(v: object, unit: str) -> str:
    x = safe_num(v)
    return f"{x} {unit}" if x != "—" else "—"


def _iter_positions(geom: dict):
    gtype = geom.get("type")
    coords = geom.get("coordinates")
    if gtype == "Point":
        yield coords
    elif gtype in ("MultiPoint", "LineString"):
        for p in coords or []:
            yield p
    elif gtype in ("MultiLineString", "Polygon"):
        for ring in coords or []:
            for p in ring:
                yield p
    elif gtype == "MultiPolygon":
        for poly in coords or []:
            for ring in poly:
                for p in ring:
                    yield p


def _bbox_from_geom(geom: dict):
    xs = []
    ys = []
    for p in _iter_positions(geom or {}):
        if not isinstance(p, (list, tuple)) or len(p) < 2:
            continue
        x = float(p[0])
        y = float(p[1])
        xs.append(x)
        ys.append(y)
    if not xs:
        return None
    return (min(xs), min(ys), max(xs), max(ys))


def _project(lon: float, lat: float, bbox: tuple[float, float, float, float], w: int, h: int):
    minx, miny, maxx, maxy = bbox
    dx = max(maxx - minx, 1e-9)
    dy = max(maxy - miny, 1e-9)
    x = (lon - minx) / dx * (w - 1)
    y = (maxy - lat) / dy * (h - 1)
    return (x, y)


def _lonlat_to_world_px(lon: float, lat: float, zoom: int) -> tuple[float, float]:
    lat = max(min(lat, 85.05112878), -85.05112878)
    n = 2**zoom
    x = (lon + 180.0) / 360.0 * n * 256.0
    lat_rad = math.radians(lat)
    y = (
        (1.0 - math.log(math.tan(lat_rad) + (1.0 / math.cos(lat_rad))) / math.pi)
        / 2.0
        * n
        * 256.0
    )
    return x, y


def _fetch_satellite_tile(z: int, x: int, y: int) -> Image.Image | None:
    if x < 0 or y < 0 or x >= 2**z or y >= 2**z:
        return None
    url = f"https://server.arcgisonline.com/ArcGIS/rest/services/World_Imagery/MapServer/tile/{z}/{y}/{x}"
    try:
        with urllib.request.urlopen(url, timeout=8) as resp:
            raw = resp.read()
        return Image.open(io.BytesIO(raw)).convert("RGB")
    except Exception:
        return None


def _build_satellite_background(
    bbox: tuple[float, float, float, float], w: int, h: int, zoom: int = 16
) -> Image.Image:
    minx, miny, maxx, maxy = bbox
    px_left, py_top = _lonlat_to_world_px(minx, maxy, zoom)
    px_right, py_bottom = _lonlat_to_world_px(maxx, miny, zoom)

    t_left = int(math.floor(px_left / 256.0))
    t_right = int(math.floor(px_right / 256.0))
    t_top = int(math.floor(py_top / 256.0))
    t_bottom = int(math.floor(py_bottom / 256.0))

    tiles_w = max(1, t_right - t_left + 1)
    tiles_h = max(1, t_bottom - t_top + 1)
    mosaic = Image.new("RGB", (tiles_w * 256, tiles_h * 256), (229, 231, 235))

    for ty in range(t_top, t_bottom + 1):
        for tx in range(t_left, t_right + 1):
            tile = _fetch_satellite_tile(zoom, tx, ty)
            if tile is None:
                continue
            mosaic.paste(tile, ((tx - t_left) * 256, (ty - t_top) * 256))

    crop_left = int(px_left - t_left * 256)
    crop_top = int(py_top - t_top * 256)
    crop_right = int(px_right - t_left * 256)
    crop_bottom = int(py_bottom - t_top * 256)
    crop_right = max(crop_right, crop_left + 1)
    crop_bottom = max(crop_bottom, crop_top + 1)

    cropped = mosaic.crop((crop_left, crop_top, crop_right, crop_bottom))
    return cropped.resize((w, h), Image.Resampling.BILINEAR)


def _draw_feature_collection(
    draw: ImageDraw.ImageDraw,
    fc: dict,
    bbox: tuple[float, float, float, float],
    w: int,
    h: int,
    *,
    fill: tuple[int, int, int, int] | None = None,
    outline: tuple[int, int, int, int] | None = None,
    width: int = 2,
):
    for f in (fc or {}).get("features", []) or []:
        geom = f.get("geometry") or {}
        gtype = geom.get("type")
        coords = geom.get("coordinates") or []
        if gtype == "Polygon":
            polys = [coords]
        elif gtype == "MultiPolygon":
            polys = coords
        else:
            polys = []
        for poly in polys:
            if not poly:
                continue
            outer = [ _project(float(p[0]), float(p[1]), bbox, w, h) for p in poly[0] if len(p) >= 2 ]
            if len(outer) >= 3:
                if fill is not None:
                    draw.polygon(outer, fill=fill)
                if outline is not None:
                    # PIL polygon outline often remains 1px; draw explicit closed line for thick borders.
                    closed = outer + [outer[0]]
                    draw.line(closed, fill=outline, width=max(1, int(width)))
        if gtype == "LineString":
            line = [_project(float(p[0]), float(p[1]), bbox, w, h) for p in coords if len(p) >= 2]
            if len(line) >= 2:
                draw.line(line, fill=outline or (249, 115, 22, 255), width=width)
        elif gtype == "MultiLineString":
            for seg in coords:
                line = [_project(float(p[0]), float(p[1]), bbox, w, h) for p in seg if len(p) >= 2]
                if len(line) >= 2:
                    draw.line(line, fill=outline or (249, 115, 22, 255), width=width)


def build_map_preview_png(payload: dict, out_png: Path) -> bool:
    overlays = payload.get("mapOverlays") or {}
    parcel = overlays.get("parcel") or payload.get("cadastreFeature")
    extent = overlays.get("extentBox")
    if not parcel or not parcel.get("geometry"):
        return False

    bbox = _bbox_from_geom((extent or {}).get("geometry") or parcel.get("geometry"))
    if not bbox:
        return False
    minx, miny, maxx, maxy = bbox
    cx = (minx + maxx) / 2
    cy = (miny + maxy) / 2
    # Охват оставляем без расширения (не меняем экстент/контекст).
    half_w = max((maxx - minx) / 12, 1e-6)
    half_h = max((maxy - miny) / 12, 1e-6)
    zoom_bbox = (cx - half_w, cy - half_h, cx + half_w, cy + half_h)

    # Более высокое разрешение для вставки в PPTX.
    w, h = 2048, 1152
    satellite = _build_satellite_background(zoom_bbox, w, h, zoom=16)
    img = satellite.convert("RGBA")
    draw = ImageDraw.Draw(img, "RGBA")
    overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    overlay_draw = ImageDraw.Draw(overlay, "RGBA")

    # Полигональные слои: сильно прозрачные, чтобы не «забивать» спутник.
    _draw_feature_collection(overlay_draw, overlays.get("waterBuffers"), zoom_bbox, w, h, fill=(96, 165, 250, 40), outline=(37, 99, 235, 210), width=5)
    _draw_feature_collection(overlay_draw, overlays.get("waterSave"), zoom_bbox, w, h, fill=(103, 232, 249, 40), outline=(8, 145, 178, 210), width=5)
    _draw_feature_collection(overlay_draw, overlays.get("powerBuffers"), zoom_bbox, w, h, fill=(239, 68, 68, 40), outline=(220, 38, 38, 210), width=5)
    _draw_feature_collection(overlay_draw, overlays.get("ooptAreas"), zoom_bbox, w, h, fill=(192, 132, 252, 40), outline=(147, 51, 234, 210), width=5)
    _draw_feature_collection(overlay_draw, overlays.get("landuseIntersected"), zoom_bbox, w, h, fill=(253, 224, 71, 40), outline=(202, 138, 4, 210), width=5)
    img = Image.alpha_composite(img, overlay)
    _draw_feature_collection(draw, overlays.get("powerLines"), zoom_bbox, w, h, outline=(249, 115, 22, 220), width=3)

    parcel_fc = {"type": "FeatureCollection", "features": [parcel]}
    # Контур участка рисуем двойным (белая подложка + красный), чтобы был максимально заметен.
    _draw_feature_collection(draw, parcel_fc, zoom_bbox, w, h, fill=None, outline=(0, 0, 0, 255), width=30)
    _draw_feature_collection(draw, parcel_fc, zoom_bbox, w, h, fill=None, outline=(255, 255, 255, 255), width=22)
    _draw_feature_collection(draw, parcel_fc, zoom_bbox, w, h, fill=None, outline=(220, 38, 38, 255), width=14)

    # Служебная метка рендера, чтобы исключить визуально «старый кэш» в отчете.
    stamp = f"render:{datetime.now().strftime('%Y-%m-%d %H:%M:%S')} v6"
    pad = 8
    text_w = max(220, len(stamp) * 7)
    text_h = 24
    x0, y0 = 12, 12
    x1, y1 = x0 + text_w + pad * 2, y0 + text_h
    draw.rectangle([x0, y0, x1, y1], fill=(0, 0, 0, 160))
    draw.text((x0 + pad, y0 + 6), stamp, fill=(255, 255, 255, 230))

    out_png.parent.mkdir(parents=True, exist_ok=True)
    img.convert("RGB").save(str(out_png), format="PNG")
    return True


def make_render_data(payload: dict, sample_data: dict, template_root: Path) -> dict:
    data = copy.deepcopy(sample_data)
    text = data.setdefault("text", {})
    images = data.setdefault("images", {})
    conditions = data.setdefault("conditions", {})

    feature = payload.get("cadastreFeature") or {}
    props = feature.get("properties") or {}
    opts = props.get("options") or {}
    terrain = payload.get("terrain") or {}
    risk_summary = [str(x) for x in (payload.get("riskSummary") or []) if str(x).strip()]
    cad_num = str_or_default(payload.get("cadastreNumber") or opts.get("cad_num") or props.get("label"))
    address = clean_address(opts.get("readable_address") or props.get("descr"))
    area_m2 = opts.get("specified_area") or opts.get("land_record_area")
    area_text = format_sotok(area_m2) if area_m2 is not None else "—"
    slope = terrain.get("maxSlopeDeg")
    elevation = terrain.get("elevationM")

    risk_score = min(100, max(0, len(risk_summary) * 30))
    risk_text, risk_img = risk_level(risk_score)
    terrain_text, terrain_img = terrain_level_name(float(slope) if slope is not None else None)

    text["kadastrnum"] = cad_num
    text["adress"] = address
    text["squire"] = area_text
    text["date"] = format_check_date_ru(datetime.now())
    text["orderer"] = "GeoRisk"
    text["riskball"] = str(risk_score)
    text["risktext"] = risk_text
    text["textaboutrisks"] = ", ".join(risk_summary) if risk_summary else "По выбранным слоям пересечений не найдено."
    text["risk1"] = risk_summary[0] if len(risk_summary) > 0 else "Нет данных"
    text["risk1dist"] = "—"
    text["risk1rule"] = "Проверить ограничения и режимы использования территории."
    text["risk2"] = risk_summary[1] if len(risk_summary) > 1 else ""
    text["risk2dist"] = ""
    text["risk2rule"] = ""
    text["risk3"] = risk_summary[2] if len(risk_summary) > 2 else ""
    text["risk3dist"] = ""
    text["risk3rule"] = ""
    text["risk4"] = risk_summary[3] if len(risk_summary) > 3 else ""
    text["risk4dist"] = ""
    text["risk4rule"] = ""
    text["releftxt"] = f"Максимальный уклон: {with_units(slope, '°')}, высота: {with_units(elevation, 'м')}."
    text["levelrelef"] = terrain_text

    images["img_sales_chart"] = str(template_root / "assets" / risk_img)
    images["levelrelefimg"] = str(template_root / "assets" / terrain_img)
    images["levellepimg"] = str(template_root / "assets" / "high.png")
    images["levelgenplanimg"] = str(template_root / "assets" / "medium.png")
    images["levelvodoimg"] = str(template_root / "assets" / "zero.png")

    conditions["block_risks"] = True
    return data


def enforce_report_text_styles(
    output_path: Path,
    *,
    address_value: str,
    cadastre_value: str,
    area_value: str,
    orderer_value: str,
    date_value: str,
) -> None:
    prs = Presentation(str(output_path))
    changed = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    txt = run.text or ""
                    if not txt.strip():
                        continue
                    # Адрес всегда обычный
                    if address_value and address_value != "—" and address_value in txt:
                        if run.font.bold is not False:
                            run.font.bold = False
                            changed = True
                    # Кадастровый номер / площадь / заказчик всегда жирные
                    if cadastre_value and cadastre_value in txt:
                        if run.font.bold is not True:
                            run.font.bold = True
                            changed = True
                    if area_value and area_value != "—" and area_value in txt:
                        if run.font.bold is not True:
                            run.font.bold = True
                            changed = True
                    if orderer_value and orderer_value in txt:
                        if run.font.bold is not True:
                            run.font.bold = True
                            changed = True
                    if date_value and date_value in txt:
                        if run.font.bold is not True:
                            run.font.bold = True
                            changed = True
    if changed:
        prs.save(str(output_path))


def main() -> int:
    args = parse_args()
    try:
        raw = sys.stdin.read()
        payload = json.loads(raw) if raw.strip() else {}
        template_root = Path(__file__).resolve().parent / "pptx_template"
        template_path = template_root / "report.pptx"
        sample_path = template_root / "data" / "sample_variant_b.json"

        with sample_path.open("r", encoding="utf-8") as f:
            sample_data = json.load(f)

        render_data = make_render_data(payload, sample_data, template_root)
        temp_map_path = Path(args.out).with_suffix(".map.png")
        if build_map_preview_png(payload, temp_map_path):
            render_data.setdefault("images", {})["img_sales_chart"] = str(temp_map_path)
        temp_data_path = Path(args.out).with_suffix(".json")
        temp_data_path.parent.mkdir(parents=True, exist_ok=True)
        with temp_data_path.open("w", encoding="utf-8") as f:
            json.dump(render_data, f, ensure_ascii=False)

        out_path = Path(args.out)
        build_report(template_path, temp_data_path, out_path)
        enforce_report_text_styles(
            out_path,
            address_value=render_data.get("text", {}).get("adress", ""),
            cadastre_value=render_data.get("text", {}).get("kadastrnum", ""),
            area_value=render_data.get("text", {}).get("squire", ""),
            orderer_value=render_data.get("text", {}).get("orderer", ""),
            date_value=render_data.get("text", {}).get("date", ""),
        )
        try:
            temp_data_path.unlink(missing_ok=True)
            temp_map_path.unlink(missing_ok=True)
        except Exception:
            pass
        print(json.dumps({"ok": True, "output": str(out_path)}, ensure_ascii=False))
        return 0
    except Exception as exc:
        print(json.dumps({"error": str(exc)}, ensure_ascii=False), file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
