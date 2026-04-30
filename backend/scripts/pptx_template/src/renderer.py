from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Any, Dict

from pptx import Presentation

PLACEHOLDER_PATTERN = re.compile(r"\{\{([a-zA-Z0-9_.-]+)\}\}")


class ReportBuildError(Exception):
    """Raised when report data or template processing fails."""


def load_report_data(data_path: str | Path) -> Dict[str, Dict[str, Any]]:
    path = Path(data_path)
    if not path.exists():
        raise ReportBuildError(f"Data file does not exist: {path}")

    with path.open("r", encoding="utf-8") as fh:
        payload = json.load(fh)

    if not isinstance(payload, dict):
        raise ReportBuildError("JSON root must be an object.")

    text = payload.get("text", {})
    images = payload.get("images", {})
    conditions = payload.get("conditions", {})

    for key, section in (
        ("text", text),
        ("images", images),
        ("conditions", conditions),
    ):
        if not isinstance(section, dict):
            raise ReportBuildError(f"'{key}' must be an object.")

    for cond_name, cond_value in conditions.items():
        if not isinstance(cond_value, bool):
            raise ReportBuildError(
                f"Condition '{cond_name}' must be boolean, got: {type(cond_value)}"
            )

    return {"text": text, "images": images, "conditions": conditions}


def _replace_placeholders_in_text(text: str, values: Dict[str, Any]) -> str:
    def _replacement(match: re.Match[str]) -> str:
        key = match.group(1)
        if key not in values:
            return ""
        return str(values[key])

    return PLACEHOLDER_PATTERN.sub(_replacement, text)


def _replace_text_on_slide(slide: Any, text_values: Dict[str, Any]) -> None:
    for shape in slide.shapes:
        _replace_text_recursive(shape, text_values)


def _hide_shapes_with_all_missing_placeholders_on_slide(
    slide: Any, text_values: Dict[str, Any]
) -> None:
    for shape in list(slide.shapes):
        _hide_shape_if_all_placeholders_missing_recursive(shape, text_values)


def _hide_shape_if_all_placeholders_missing_recursive(
    shape: Any, text_values: Dict[str, Any]
) -> None:
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for child in list(shape.shapes):
            _hide_shape_if_all_placeholders_missing_recursive(child, text_values)
        return

    keys = _collect_placeholder_keys_in_shape(shape)
    if keys and all((key not in text_values) for key in keys):
        element = shape._element
        parent = element.getparent()
        parent.remove(element)


def _collect_placeholder_keys_in_shape(shape: Any) -> set[str]:
    keys: set[str] = set()
    if getattr(shape, "has_text_frame", False):
        keys |= _collect_placeholder_keys_in_text_frame(shape.text_frame)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                keys |= _collect_placeholder_keys_in_text_frame(cell.text_frame)
    return keys


def _collect_placeholder_keys_in_text_frame(text_frame: Any) -> set[str]:
    keys: set[str] = set()
    for paragraph in text_frame.paragraphs:
        if paragraph.runs:
            text = "".join(run.text for run in paragraph.runs)
        else:
            text = paragraph.text
        keys.update(PLACEHOLDER_PATTERN.findall(text))
    return keys


def _replace_text_recursive(shape: Any, text_values: Dict[str, Any]) -> None:
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for child in shape.shapes:
            _replace_text_recursive(child, text_values)
        return

    if getattr(shape, "has_text_frame", False):
        _replace_text_in_text_frame(shape.text_frame, text_values)

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                _replace_text_in_text_frame(cell.text_frame, text_values)


def _replace_text_in_text_frame(text_frame: Any, text_values: Dict[str, Any]) -> None:
    for paragraph in text_frame.paragraphs:
        _replace_text_in_paragraph(paragraph, text_values)


def _replace_text_in_paragraph(paragraph: Any, text_values: Dict[str, Any]) -> None:
    if paragraph.runs:
        _replace_placeholders_preserving_runs(paragraph, text_values)
        return

    # Fallback for rare paragraphs without runs.
    paragraph.text = _replace_placeholders_in_text(paragraph.text, text_values)


def _replace_placeholders_preserving_runs(paragraph: Any, text_values: Dict[str, Any]) -> None:
    # Replace from right to left so character offsets remain valid.
    while True:
        runs = list(paragraph.runs)
        if not runs:
            return

        full_text = "".join(run.text for run in runs)
        matches = list(PLACEHOLDER_PATTERN.finditer(full_text))
        if not matches:
            return

        changed = False
        for match in reversed(matches):
            key = match.group(1)
            if key not in text_values:
                continue

            replacement = str(text_values[key])
            start, end = match.span()
            run_map = _build_run_map(runs)
            start_run_idx, start_off = _locate_offset(run_map, start)
            end_run_idx, end_off = _locate_offset(run_map, end)

            if start_run_idx == end_run_idx:
                run = runs[start_run_idx]
                run.text = run.text[:start_off] + replacement + run.text[end_off:]
            else:
                start_run = runs[start_run_idx]
                end_run = runs[end_run_idx]
                prefix = start_run.text[:start_off]
                suffix = end_run.text[end_off:]

                start_run.text = prefix + replacement
                for idx in range(start_run_idx + 1, end_run_idx):
                    runs[idx].text = ""
                end_run.text = suffix

            changed = True

        if not changed:
            return


def _build_run_map(runs: Any) -> list[tuple[int, int, int]]:
    mapping = []
    cursor = 0
    for idx, run in enumerate(runs):
        text = run.text or ""
        start = cursor
        cursor += len(text)
        mapping.append((idx, start, cursor))
    return mapping


def _locate_offset(run_map: list[tuple[int, int, int]], char_pos: int) -> tuple[int, int]:
    for idx, start, end in run_map:
        if start <= char_pos < end:
            return idx, char_pos - start
        if char_pos == end:
            return idx, end - start
    last_idx, last_start, last_end = run_map[-1]
    return last_idx, last_end - last_start


def _replace_images_on_slide(slide: Any, image_map: Dict[str, str]) -> None:
    # Iterate over a snapshot to avoid replacing newly inserted pictures again.
    for shape in list(slide.shapes):
        _replace_image_recursive(slide, shape, image_map)


def _replace_image_recursive(slide: Any, shape: Any, image_map: Dict[str, str]) -> None:
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for child in shape.shapes:
            _replace_image_recursive(slide, child, image_map)
        return

    target_path = image_map.get(shape.name)
    if not target_path:
        return

    image_path = Path(target_path)
    if not image_path.exists():
        raise ReportBuildError(
            f"Image for '{shape.name}' does not exist: {image_path}"
        )

    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height
    parent = shape._element.getparent()
    parent.remove(shape._element)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height).name = (
        shape.name
    )


def _apply_conditions_on_slide(slide: Any, conditions: Dict[str, bool]) -> None:
    for shape in list(slide.shapes):
        _apply_condition_recursive(shape, conditions)


def _apply_condition_recursive(shape: Any, conditions: Dict[str, bool]) -> None:
    # If a parent group is hidden, remove it entirely without traversing children.
    if shape.name in conditions and not conditions[shape.name]:
        element = shape._element
        parent = element.getparent()
        parent.remove(element)
        return

    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for child in list(shape.shapes):
            _apply_condition_recursive(child, conditions)


def build_report(template_path: str | Path, data_path: str | Path, output_path: str | Path) -> None:
    template = Path(template_path)
    if not template.exists():
        raise ReportBuildError(f"Template does not exist: {template}")

    data = load_report_data(data_path)
    presentation = Presentation(str(template))

    for slide in presentation.slides:
        _hide_shapes_with_all_missing_placeholders_on_slide(slide, data["text"])
        _replace_text_on_slide(slide, data["text"])
        _replace_images_on_slide(slide, data["images"])
        _apply_conditions_on_slide(slide, data["conditions"])

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output))
